"""Office documents ability — create and read Word, Excel, PowerPoint files.

Create actions:
- create_docx: Microsoft Word (.docx) from markdown-ish text
- create_xlsx: Microsoft Excel (.xlsx) from structured sheets data
- create_pptx: Microsoft PowerPoint (.pptx) from slide outline

Read actions:
- read_docx: Extract text from .docx (base64 or path)
- read_xlsx: Extract sheets/rows from .xlsx (base64 or path)
- read_pptx: Extract slides text from .pptx (base64 or path)

Also auto-extracts uploaded docx/xlsx/pptx via pre_process — LLM sees
the document content as plain text without needing to call any tool.

Dependencies (lazy-installed via ensure_dependencies):
- python-docx (Word)
- openpyxl (Excel)
- python-pptx (PowerPoint)
"""

from __future__ import annotations

import asyncio
import json as _json
import logging
import os
import re
import sys
import time
from copy import deepcopy

from syne.abilities.base import Ability

logger = logging.getLogger("syne.ability.office")


# ── Template resolution ───────────────────────────────────────────────────
_TEMPLATES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")
_DOCX_TEMPLATE = os.path.join(_TEMPLATES_DIR, "skeleton.docx")
_PPTX_TEMPLATE = os.path.join(_TEMPLATES_DIR, "skeleton.pptx")


def _docx_template_path() -> str | None:
    return _DOCX_TEMPLATE if os.path.isfile(_DOCX_TEMPLATE) else None


def _pptx_template_path() -> str | None:
    return _PPTX_TEMPLATE if os.path.isfile(_PPTX_TEMPLATE) else None


# PPTX template — 9 demo slides indexed by layout name
_PPTX_LAYOUTS = [
    "cover",         # 0
    "section",       # 1
    "two_column",    # 2
    "icon_rows",     # 3
    "stat_callouts", # 4
    "grid_2x2",      # 5
    "timeline",      # 6
    "comparison",    # 7
    "closing",       # 8
]


def _now_ts() -> str:
    return time.strftime("%Y%m%d_%H%M%S")


def _safe_filename(name: str, default: str = "document") -> str:
    name = (name or "").strip() or default
    name = re.sub(r"[^a-zA-Z0-9._-]+", "_", name)
    name = re.sub(r"_+", "_", name).strip("_")
    return name[:120] or default


def _ensure_ext(name: str, ext: str) -> str:
    n = (name or "").strip() or "output"
    ext = ext if ext.startswith(".") else f".{ext}"
    while n.lower().endswith(ext.lower()):
        n = n[: -len(ext)]
    n = n.rstrip(". ")
    return (n or "output") + ext


class OfficeAbility(Ability):
    name = "office"
    description = "Create and read Microsoft Office documents (Word, Excel, PowerPoint)."
    version = "2.0"
    permission = 0o770
    # Priority pre-processing: auto-extract content from uploaded .docx/.xlsx/.pptx
    priority = True

    _DEPS = {
        "python-docx": "docx",
        "openpyxl": "openpyxl",
        "python-pptx": "pptx",
    }

    # MIME mapping for pre_process
    _DOCX_MIMES = {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    }
    _XLSX_MIMES = {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    }
    _PPTX_MIMES = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }

    # Read limits — truncate to keep LLM context manageable
    _MAX_EXTRACTED_CHARS = 50_000
    _MAX_XLSX_ROWS_PER_SHEET = 200
    _MAX_XLSX_COLS_PER_SHEET = 50

    def handles_input_type(self, input_type: str) -> bool:
        return input_type == "document"

    async def pre_process(
        self, input_type: str, input_data: dict, user_prompt: str,
        config: dict | None = None,
    ) -> str | None:
        """Auto-extract text from uploaded Office docs. Returns None for other types
        (falls through to next handler, e.g. pdf ability)."""
        mime = (input_data.get("mime_type") or "").lower()
        filename = (input_data.get("filename") or "").lower()

        kind = None
        if mime in self._DOCX_MIMES or filename.endswith(".docx") or filename.endswith(".doc"):
            kind = "docx"
        elif mime in self._XLSX_MIMES or filename.endswith(".xlsx") or filename.endswith(".xls"):
            kind = "xlsx"
        elif mime in self._PPTX_MIMES or filename.endswith(".pptx") or filename.endswith(".ppt"):
            kind = "pptx"
        else:
            return None  # not an Office doc — let other abilities handle

        b64 = input_data.get("base64", "")
        path = input_data.get("path", "")
        if not b64 and not path:
            return None

        try:
            import base64
            if b64:
                content = base64.b64decode(b64)
            else:
                with open(path, "rb") as f:
                    content = f.read()
        except Exception as e:
            logger.warning(f"Office pre_process: read failed: {e}")
            return None

        try:
            if kind == "docx":
                text, meta = _read_docx_bytes(content)
                header = f"Word: {filename or 'document.docx'} ({meta['paragraphs']} paragraphs, {len(text)} chars)"
            elif kind == "xlsx":
                text, meta = _read_xlsx_bytes(
                    content,
                    max_rows=self._MAX_XLSX_ROWS_PER_SHEET,
                    max_cols=self._MAX_XLSX_COLS_PER_SHEET,
                )
                header = f"Excel: {filename or 'workbook.xlsx'} ({meta['sheets']} sheet(s), {len(text)} chars)"
            else:  # pptx
                text, meta = _read_pptx_bytes(content)
                header = f"PowerPoint: {filename or 'slides.pptx'} ({meta['slides']} slide(s), {len(text)} chars)"
        except Exception as e:
            logger.warning(f"Office pre_process: extraction failed ({kind}): {e}")
            return None

        if len(text) > self._MAX_EXTRACTED_CHARS:
            text = text[: self._MAX_EXTRACTED_CHARS] + f"\n\n[... truncated at {self._MAX_EXTRACTED_CHARS} chars]"

        return f"{header}\n\n{text}" if text.strip() else f"{header}\n\n[empty document]"

    async def ensure_dependencies(self) -> tuple[bool, str]:
        """Install Office deps (python-docx, openpyxl, python-pptx)."""
        missing_pkgs = []
        for pip_name, import_name in self._DEPS.items():
            try:
                __import__(import_name)
            except ImportError:
                missing_pkgs.append(pip_name)

        if not missing_pkgs:
            return True, ""

        logger.info(f"Installing Office deps: {', '.join(missing_pkgs)}")
        try:
            proc = await asyncio.create_subprocess_exec(
                sys.executable, "-m", "pip", "install", *missing_pkgs,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            _, stderr = await asyncio.wait_for(proc.communicate(), timeout=1800)
        except asyncio.TimeoutError:
            return False, f"Install timed out. Manual: pip install {' '.join(missing_pkgs)}"
        if proc.returncode != 0:
            err = stderr.decode().strip()
            return False, f"Failed to install Office deps: {err}"
        return True, f"Installed: {', '.join(missing_pkgs)}"

    def get_schema(self) -> dict:
        return {
            "type": "function",
            "function": {
                "name": "office",
                "description": (
                    "Create or read Microsoft Office documents. Create actions: "
                    "create_docx (Word), create_xlsx (Excel), create_pptx (PowerPoint). "
                    "Read actions: read_docx, read_xlsx, read_pptx (from base64). "
                    "Uploaded Office files are auto-extracted via pre_process — no tool call needed."
                ),
                "parameters": {
                    "type": "object",
                    "properties": {
                        "action": {
                            "type": "string",
                            "enum": [
                                "create_docx", "create_xlsx", "create_pptx",
                                "read_docx", "read_xlsx", "read_pptx",
                            ],
                        },
                        "title": {
                            "type": "string",
                            "description": "Document title (used as filename if out_name not given).",
                        },
                        "content": {
                            "type": "string",
                            "description": (
                                "For create_docx: body text. Markdown-style: "
                                "# H1, ## H2, ### H3, - bullets, **bold**, *italic*. "
                                "Paragraphs separated by blank line."
                            ),
                        },
                        "sheets": {
                            "type": "string",
                            "description": (
                                "For create_xlsx: JSON array of sheets. Each sheet: "
                                '{"name": "Sheet1", "headers": ["A","B"], "rows": [[1,2],[3,4]]}. '
                                "Headers optional. Pass as JSON string."
                            ),
                        },
                        "slides": {
                            "type": "string",
                            "description": (
                                "For create_pptx: JSON array of slides. The skeleton template provides 9 "
                                "layouts and Syne auto-picks the best one per slide. Layout is chosen by "
                                "which fields you provide:\n"
                                "- cover (auto, first slide): {title, subtitle, meta}\n"
                                "- section: {title, subtitle} (no bullets/content)\n"
                                "- stat_callouts: {title, subtitle, stats: [{value,label}, ...]} (max 3)\n"
                                "- timeline: {title, subtitle, phases: [{name,desc}, ...]} (max 5)\n"
                                "- comparison: {title, subtitle, left: {name,bullets:[]}, right: {name,bullets:[],recommended:true}}\n"
                                "- grid_2x2: {title, subtitle, items: [{name,desc}, x4]}\n"
                                "- icon_rows: {title, subtitle, items: [{name,desc}, ...]} (not exactly 4)\n"
                                "- two_column (default): {title, subtitle, left, right} or {title, bullets}\n"
                                "- closing (last slide, auto): {title:'Terima Kasih'|'Thank You'|..., subtitle, meta}\n"
                                "Pass as JSON string."
                            ),
                        },
                        "out_name": {"type": "string", "description": "Optional output filename."},
                        "file_base64": {
                            "type": "string",
                            "description": "For read_* actions: base64-encoded file bytes.",
                        },
                    },
                    "required": ["action"],
                },
            },
        }

    def get_guide(self, enabled: bool, config: dict) -> str:
        if not enabled:
            return (
                "- Status: **disabled**\n"
                "- Enable: `update_ability(action='enable', name='office')`"
            )

        missing = []
        for mod in self._DEPS.values():
            try:
                __import__(mod)
            except Exception:
                missing.append(mod)

        if missing:
            return (
                f"- Status: **not ready** (missing deps: {', '.join(missing)})\n"
                "- Will auto-install on first use, or run `pip install python-docx openpyxl python-pptx`"
            )

        return (
            "- Status: **ready**\n"
            "- Auto-extracts content from uploaded .docx/.xlsx/.pptx (priority pre-process)\n"
            "- Create Word: `office(action='create_docx', title='...', content='# Heading\\nText')`\n"
            "- Create Excel: `office(action='create_xlsx', title='...', sheets='[{\"name\":\"S1\",\"headers\":[\"A\"],\"rows\":[[1]]}]')`\n"
            "- Create PPT: `office(action='create_pptx', title='...', slides='[{\"title\":\"S1\",\"bullets\":[\"p1\"]}]')`\n"
            "- Read base64: `office(action='read_docx'|'read_xlsx'|'read_pptx', file_base64='...')`"
        )

    async def execute(self, params: dict, context: dict) -> dict:
        action = (params.get("action") or "").strip()
        title = params.get("title") or ""
        out_name = params.get("out_name") or ""

        session_id = str(context.get("session_id") or "")
        out_dir = self.get_output_dir(session_id=session_id)
        os.makedirs(out_dir, exist_ok=True)

        try:
            if action == "create_docx":
                content = params.get("content") or ""
                name = _ensure_ext(_safe_filename(out_name or title or f"document_{_now_ts()}"), ".docx")
                out_path = os.path.join(out_dir, name)
                _make_docx(out_path, title, content)
                return {"success": True, "result": {"file_path": out_path}, "media": out_path}

            if action == "create_xlsx":
                sheets_str = params.get("sheets") or ""
                if not sheets_str:
                    raise ValueError("sheets is required (JSON array)")
                try:
                    sheets = _json.loads(sheets_str) if isinstance(sheets_str, str) else sheets_str
                except _json.JSONDecodeError as e:
                    raise ValueError(f"Invalid sheets JSON: {e}")
                if not isinstance(sheets, list) or not sheets:
                    raise ValueError("sheets must be a non-empty JSON array")
                name = _ensure_ext(_safe_filename(out_name or title or f"workbook_{_now_ts()}"), ".xlsx")
                out_path = os.path.join(out_dir, name)
                _make_xlsx(out_path, sheets)
                return {"success": True, "result": {"file_path": out_path}, "media": out_path}

            if action == "create_pptx":
                slides_str = params.get("slides") or ""
                if not slides_str:
                    raise ValueError("slides is required (JSON array)")
                try:
                    slides = _json.loads(slides_str) if isinstance(slides_str, str) else slides_str
                except _json.JSONDecodeError as e:
                    raise ValueError(f"Invalid slides JSON: {e}")
                if not isinstance(slides, list) or not slides:
                    raise ValueError("slides must be a non-empty JSON array")
                name = _ensure_ext(_safe_filename(out_name or title or f"slides_{_now_ts()}"), ".pptx")
                out_path = os.path.join(out_dir, name)
                _make_pptx(out_path, title, slides)
                return {"success": True, "result": {"file_path": out_path}, "media": out_path}

            if action in ("read_docx", "read_xlsx", "read_pptx"):
                b64 = params.get("file_base64") or ""
                if not b64:
                    raise ValueError("file_base64 is required")
                import base64 as _b64
                content = _b64.b64decode(b64)
                if action == "read_docx":
                    text, meta = _read_docx_bytes(content)
                elif action == "read_xlsx":
                    text, meta = _read_xlsx_bytes(
                        content,
                        max_rows=self._MAX_XLSX_ROWS_PER_SHEET,
                        max_cols=self._MAX_XLSX_COLS_PER_SHEET,
                    )
                else:  # read_pptx
                    text, meta = _read_pptx_bytes(content)
                if len(text) > self._MAX_EXTRACTED_CHARS:
                    text = text[: self._MAX_EXTRACTED_CHARS]
                    meta["truncated"] = True
                return {"success": True, "result": {**meta, "text": text}}

            raise ValueError(f"Unknown action: {action}")

        except Exception as e:
            return {"success": False, "error": str(e)}


# ─────────────────────────────────────────────────────────────────────────
# Readers
# ─────────────────────────────────────────────────────────────────────────


def _read_docx_bytes(content: bytes) -> tuple[str, dict]:
    """Extract text from .docx bytes. Returns (text, meta)."""
    import io
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts = []
    para_count = 0
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        # Convert heading style to markdown-ish
        if "heading" in style:
            level_match = re.search(r"(\d+)", style)
            level = int(level_match.group(1)) if level_match else 1
            parts.append(("#" * min(level, 6)) + " " + t)
        else:
            parts.append(t)
        para_count += 1

    # Also extract tables
    table_count = 0
    for tbl in doc.tables:
        table_count += 1
        parts.append(f"\n[Table {table_count}]")
        for row in tbl.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))

    text = "\n\n".join(parts)
    return text, {"paragraphs": para_count, "tables": table_count}


def _read_xlsx_bytes(content: bytes, max_rows: int = 200, max_cols: int = 50) -> tuple[str, dict]:
    """Extract sheets/rows from .xlsx bytes. Returns (text, meta).

    Each sheet is rendered as Markdown-style table preview.
    """
    import io
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    sheet_count = 0
    total_rows = 0
    for ws in wb.worksheets:
        sheet_count += 1
        parts.append(f"--- Sheet: {ws.title} ---")
        rows_emitted = 0
        truncated_rows = False
        truncated_cols = False
        for row in ws.iter_rows(values_only=True):
            if rows_emitted >= max_rows:
                truncated_rows = True
                break
            cells = list(row)
            if len(cells) > max_cols:
                cells = cells[:max_cols]
                truncated_cols = True
            # Skip fully empty rows
            if all(c is None or (isinstance(c, str) and not c.strip()) for c in cells):
                continue
            parts.append(" | ".join("" if c is None else str(c) for c in cells))
            rows_emitted += 1
            total_rows += 1
        if truncated_rows:
            parts.append(f"[... row limit {max_rows} reached]")
        if truncated_cols:
            parts.append(f"[... col limit {max_cols} applied]")
        parts.append("")  # blank line between sheets

    wb.close()
    text = "\n".join(parts).rstrip()
    return text, {"sheets": sheet_count, "rows": total_rows}


def _read_pptx_bytes(content: bytes) -> tuple[str, dict]:
    """Extract slides text from .pptx bytes. Returns (text, meta)."""
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts = []
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
        slide_lines = [f"--- Slide {slide_count} ---"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip()
                if t:
                    slide_lines.append(t)
        parts.append("\n".join(slide_lines))

    text = "\n\n".join(parts)
    return text, {"slides": slide_count}



# ── Arabic / RTL helpers ──────────────────────────────────────────────────
_ARABIC_RE = re.compile(r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]")

def _has_arabic(text: str) -> bool:
    """Return True if text contains Arabic characters."""
    return bool(_ARABIC_RE.search(text or ""))

def _set_run_arabic(run) -> None:
    """Set font to Arial (Latin + Complex Script) and enable RTL on a Run containing Arabic."""
    from docx.oxml.ns import qn
    run.font.name = "Arial"
    # Set Complex Script font (w:cs on w:rFonts) — this is what Word uses for Arabic glyphs
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = run._r.makeelement(qn("w:rFonts"), {})
        rPr.insert(0, rFonts)
    rFonts.set(qn("w:cs"), "Arial")
    # Enable RTL on run properties
    rtl = rPr.find(qn("w:rtl"))
    if rtl is None:
        rtl = run._r.makeelement(qn("w:rtl"), {})
        rPr.append(rtl)
    rtl.set(qn("w:val"), "1")
    # Enable Complex Script explicitly
    cs = rPr.find(qn("w:cs"))
    if cs is None:
        cs = run._r.makeelement(qn("w:cs"), {})
        rPr.append(cs)
    cs.set(qn("w:val"), "1")

def _set_paragraph_bidi(paragraph) -> None:
    """Enable bidi on paragraph properties for proper RTL rendering."""
    from docx.oxml.ns import qn
    pPr = paragraph._p.get_or_add_pPr()
    bidi = pPr.find(qn("w:bidi"))
    if bidi is None:
        bidi = paragraph._p.makeelement(qn("w:bidi"), {})
        pPr.append(bidi)
    bidi.set(qn("w:val"), "1")

# ─────────────────────────────────────────────────────────────────────────
# Builders
# ─────────────────────────────────────────────────────────────────────────


def _add_hyperlink(paragraph, label: str, url: str) -> None:
    """Append a clickable hyperlink run to a python-docx paragraph."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    part = paragraph.part
    r_id = part.relate_to(
        url,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)

    new_run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "1A73E8")
    rPr.append(color)
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    rPr.append(underline)
    new_run.append(rPr)
    t = OxmlElement("w:t")
    t.text = label
    t.set(qn("xml:space"), "preserve")
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


# Regex patterns reused across the inline tokenizer.
_INLINE_CODE_RE = re.compile(r"`([^`\n]+)`")
_INLINE_LINK_RE = re.compile(r"\[([^\]\n]+)\]\(([^)\n\s]+)\)")
_INLINE_BOLD_RE = re.compile(r"\*\*([^*\n]+)\*\*")
_INLINE_STRIKE_RE = re.compile(r"~~([^~\n]+)~~")
_INLINE_ITALIC_RE = re.compile(r"(?<![\*\w])\*([^*\n]+?)\*(?![\*\w])")
_INLINE_ITALIC_UNDER_RE = re.compile(r"(?<![_\w])_([^_\n]+?)_(?![_\w])")


def _tokenize_inline_md(text: str) -> list[tuple[str, object]]:
    """Tokenize a string into a flat list of styled segments.

    Returns: list of (kind, payload) where kind is one of
    'text', 'bold', 'italic', 'code', 'strike', 'link'.
    For 'link' payload is (label, url); for others it is a string.
    """
    tokens: list[tuple[str, object]] = [("text", text)]

    def _split(kind: str, pattern: re.Pattern, extract):
        out: list[tuple[str, object]] = []
        for tk, payload in tokens:
            if tk != "text":
                out.append((tk, payload))
                continue
            assert isinstance(payload, str)
            last = 0
            for m in pattern.finditer(payload):
                if m.start() > last:
                    out.append(("text", payload[last:m.start()]))
                out.append((kind, extract(m)))
                last = m.end()
            if last < len(payload):
                out.append(("text", payload[last:]))
        tokens[:] = out

    _split("code", _INLINE_CODE_RE, lambda m: m.group(1))
    _split("link", _INLINE_LINK_RE, lambda m: (m.group(1), m.group(2)))
    _split("bold", _INLINE_BOLD_RE, lambda m: m.group(1))
    _split("strike", _INLINE_STRIKE_RE, lambda m: m.group(1))
    _split("italic", _INLINE_ITALIC_RE, lambda m: m.group(1))
    _split("italic", _INLINE_ITALIC_UNDER_RE, lambda m: m.group(1))
    return [t for t in tokens if t[1] != ""]


def _add_inline_runs(paragraph, text: str) -> None:
    """Add text to a docx paragraph honoring markdown inline markers.

    Supports **bold**, *italic*, _italic_, `code`, ~~strike~~, [text](url).
    Arabic text gets Arial font + RTL direction automatically.
    """
    has_any_arabic = _has_arabic(text)
    for kind, payload in _tokenize_inline_md(text):
        if kind == "link":
            label, url = payload  # type: ignore[misc]
            try:
                _add_hyperlink(paragraph, label, url)
            except Exception:
                # Fallback: render as plain "label (url)"
                run = paragraph.add_run(f"{label} ({url})")
                if _has_arabic(label):
                    _set_run_arabic(run)
            continue
        assert isinstance(payload, str)
        run = paragraph.add_run(payload)
        if kind == "bold":
            run.bold = True
        elif kind == "italic":
            run.italic = True
        elif kind == "strike":
            run.font.strike = True
        elif kind == "code":
            run.font.name = "Consolas"
        if _has_arabic(payload):
            _set_run_arabic(run)

    # If any Arabic was found, set bidi on the paragraph
    if has_any_arabic:
        _set_paragraph_bidi(paragraph)


def _clear_docx_body(doc) -> None:
    """Remove all body paragraphs and tables, keep sectPr (page setup),
    headers, footers, and styles."""
    from docx.oxml.ns import qn
    body = doc.element.body
    for child in list(body):
        # Keep section properties (final element with page layout)
        if child.tag == qn("w:sectPr"):
            continue
        body.remove(child)


def _docx_body_style(doc) -> str:
    """Pick the best body paragraph style available in the template."""
    style_names = {s.name for s in doc.styles}
    for candidate in ("Body", "Body Text", "Normal"):
        if candidate in style_names:
            return candidate
    return "Normal"


def _docx_resolve_style(doc, *candidates: str) -> str | None:
    """Find the first usable style name from candidates.

    Templates can have quirky styleId vs name combos (e.g. styleId='Heading1'
    but display name 'Heading 1'). Try both forms.
    """
    for cand in candidates:
        # Try direct
        try:
            doc.styles[cand]
            return cand
        except KeyError:
            pass
        # Try without spaces
        nospace = cand.replace(" ", "")
        if nospace != cand:
            try:
                doc.styles[nospace]
                return nospace
            except KeyError:
                pass
    return None


def _docx_add_paragraph(doc, text: str, *style_candidates: str):
    """Add a paragraph using the first available style from candidates."""
    style = _docx_resolve_style(doc, *style_candidates)
    if style:
        try:
            return doc.add_paragraph(text, style=style)
        except KeyError:
            pass
    return doc.add_paragraph(text)


def _is_pipe_separator_line(line: str) -> bool:
    """Match a markdown table separator row like |---|:--:|---:|."""
    s = line.strip()
    if "|" not in s:
        return False
    core = s.strip("|")
    cells = [c.strip() for c in core.split("|")]
    if not cells:
        return False
    return all(re.fullmatch(r":?-{1,}:?", c or "") for c in cells)


def _parse_table_alignments(separator_line: str) -> list[str]:
    """Return per-column alignment ('left'|'center'|'right') from the separator row."""
    s = separator_line.strip().strip("|")
    out: list[str] = []
    for c in s.split("|"):
        c = c.strip()
        if c.startswith(":") and c.endswith(":"):
            out.append("center")
        elif c.endswith(":"):
            out.append("right")
        else:
            out.append("left")
    return out


def _normalize_md_blocks(content: str) -> str:
    """Pre-pass that protects multi-line constructs from a naive split('\\n\\n').

    1. Fenced ``` code blocks: replace their internal newlines with the sentinel
       \\x01 so the whole fence becomes a single 'block' after splitting; callers
       detect \\x01 and unjoin.
    2. Pipe tables: insert blank lines before/after a table run so the block
       boundary aligns with the table boundary, even when the LLM forgets to
       leave a blank line around the table.
    """
    lines = content.splitlines()

    # Pass 1: collapse fenced code blocks into a single line with \x01 separators.
    collapsed: list[str] = []
    fence_buf: list[str] | None = None
    for ln in lines:
        is_fence_marker = ln.strip().startswith("```")
        if fence_buf is None and is_fence_marker:
            fence_buf = [ln]
        elif fence_buf is not None and is_fence_marker:
            fence_buf.append(ln)
            collapsed.append("\x01".join(fence_buf))
            fence_buf = None
        elif fence_buf is not None:
            fence_buf.append(ln)
        else:
            collapsed.append(ln)
    if fence_buf is not None:
        # Unclosed fence — emit as-is so the user at least sees the content.
        collapsed.append("\x01".join(fence_buf))

    # Pass 2: classify each line as belonging to a special block (table or
    # fenced code) so we can insert blank lines around them. Fenced code is
    # already a single \x01-joined line at this point.
    n = len(collapsed)
    in_special = [False] * n
    i = 0
    while i < n:
        ln = collapsed[i]
        if "\x01" in ln:  # collapsed fence — its own block
            in_special[i] = True
            i += 1
            continue
        s = ln.strip()
        if s.startswith("|") and "|" in s and i + 1 < n and "\x01" not in collapsed[i + 1]:
            if _is_pipe_separator_line(collapsed[i + 1].strip()):
                in_special[i] = True
                in_special[i + 1] = True
                j = i + 2
                while j < n and "\x01" not in collapsed[j] and collapsed[j].strip().startswith("|"):
                    in_special[j] = True
                    j += 1
                i = j
                continue
        i += 1

    out: list[str] = []
    prev_special = False
    for idx, ln in enumerate(collapsed):
        cur = in_special[idx]
        # Insert a blank line on any transition (in→out or out→in) to break
        # the block boundary, so split('\n\n') isolates fences and tables.
        if cur != prev_special and out and out[-1].strip():
            out.append("")
        out.append(ln)
        prev_special = cur
    return "\n".join(out)


def _is_pipe_table(block: str) -> bool:
    """Detect a GitHub-style pipe table block.

    Requires at least 2 lines where every non-empty line contains a pipe,
    and the 2nd line is a separator row (e.g. |---|:--:|---|).
    """
    lines = [ln for ln in block.splitlines() if ln.strip()]
    if len(lines) < 2:
        return False
    if not all("|" in ln for ln in lines):
        return False
    sep = lines[1].strip().strip("|")
    cells = [c.strip() for c in sep.split("|")]
    return bool(cells) and all(re.fullmatch(r":?-{1,}:?", c) for c in cells)

def _split_pipe_row(line: str) -> list[str]:
    r"""Split a markdown table row into cell strings, honoring escaped \|."""
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    # Split on unescaped pipes
    cells = re.split(r"(?<!\\)\|", s)
    return [c.replace("\\|", "|").strip() for c in cells]

def _set_table_borders(table) -> None:
    """Force a visible single-line grid on a Word table via explicit OOXML.

    Templates may lack a 'Table Grid' style, leaving tables borderless.
    Injecting w:tblBorders guarantees gridlines regardless of style.
    """
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tblPr = table._tbl.tblPr
    existing = tblPr.find(qn("w:tblBorders"))
    if existing is not None:
        tblPr.remove(existing)
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), "4")
        el.set(qn("w:space"), "0")
        el.set(qn("w:color"), "888888")
        borders.append(el)
    tblPr.append(borders)

def _set_cell_shading(cell, hex_color: str) -> None:
    """Apply a solid background fill to a single table cell via OOXML w:shd."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    tcPr = cell._tc.get_or_add_tcPr()
    existing = tcPr.find(qn("w:shd"))
    if existing is not None:
        tcPr.remove(existing)
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)

def _docx_apply_alignment(paragraph, alignment: str) -> None:
    """Apply 'left' | 'center' | 'right' alignment to a paragraph."""
    try:
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    except ImportError:
        return
    if alignment == "center":
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    elif alignment == "right":
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.RIGHT
    else:
        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT


def _add_docx_table(doc, block: str) -> None:
    """Render a markdown pipe-table block as a native Word table.

    Column alignment is read from the separator row (:--- / :---: / ---:).
    """
    lines = [ln for ln in block.splitlines() if ln.strip()]
    header = _split_pipe_row(lines[0])
    aligns = _parse_table_alignments(lines[1]) if len(lines) > 1 else []
    body_rows = [_split_pipe_row(ln) for ln in lines[2:]]
    ncols = max([len(header)] + [len(r) for r in body_rows]) if body_rows else len(header)
    # Pad alignments to ncols
    while len(aligns) < ncols:
        aligns.append("left")

    style = _docx_resolve_style(
        doc, "Table Grid", "Light Grid", "Light List Accent 1", "Normal Table"
    )
    try:
        table = doc.add_table(rows=1, cols=ncols, style=style) if style else doc.add_table(rows=1, cols=ncols)
    except KeyError:
        table = doc.add_table(rows=1, cols=ncols)
    try:
        table.autofit = True
    except Exception:
        pass

    _set_table_borders(table)

    # Header row (bold + shaded + aligned)
    hdr_cells = table.rows[0].cells
    for i in range(ncols):
        text = header[i] if i < len(header) else ""
        cell = hdr_cells[i]
        _set_cell_shading(cell, "D9D9D9")
        para = cell.paragraphs[0]
        _add_inline_runs(para, text)
        for run in para.runs:
            run.bold = True
        _docx_apply_alignment(para, aligns[i])

    # Body rows
    for row in body_rows:
        cells = table.add_row().cells
        for i in range(ncols):
            text = row[i] if i < len(row) else ""
            para = cells[i].paragraphs[0]
            _add_inline_runs(para, text)
            _docx_apply_alignment(para, aligns[i])


def _add_docx_code_block(doc, code_text: str) -> None:
    """Render a fenced code block as monospace lines with light-gray shading."""
    from docx.shared import Pt
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    code_style = _docx_resolve_style(doc, "HTML Preformatted", "No Spacing", "Normal")

    for line in code_text.split("\n"):
        if code_style:
            try:
                p = doc.add_paragraph(style=code_style)
            except KeyError:
                p = doc.add_paragraph()
        else:
            p = doc.add_paragraph()

        pPr = p._p.get_or_add_pPr()
        # Paragraph shading (light gray)
        existing_shd = pPr.find(qn("w:shd"))
        if existing_shd is not None:
            pPr.remove(existing_shd)
        shd = OxmlElement("w:shd")
        shd.set(qn("w:val"), "clear")
        shd.set(qn("w:color"), "auto")
        shd.set(qn("w:fill"), "F3F3F3")
        pPr.append(shd)
        # Tight spacing so multiple code lines stay together
        existing_sp = pPr.find(qn("w:spacing"))
        if existing_sp is not None:
            pPr.remove(existing_sp)
        sp = OxmlElement("w:spacing")
        sp.set(qn("w:after"), "0")
        sp.set(qn("w:before"), "0")
        sp.set(qn("w:line"), "260")
        sp.set(qn("w:lineRule"), "auto")
        pPr.append(sp)

        # Empty line gets a single space so the shading still renders.
        run = p.add_run(line if line else " ")
        run.font.name = "Consolas"
        run.font.size = Pt(10)
        # Preserve leading/trailing whitespace inside the run
        t = run._r.find(qn("w:t"))
        if t is not None:
            t.set(qn("xml:space"), "preserve")

def _make_docx(out_path: str, title: str, content: str) -> None:
    from docx import Document
    from docx.shared import Pt

    template = _docx_template_path()
    if template:
        doc = Document(template)
        _clear_docx_body(doc)
    else:
        doc = Document()

    # Title — try Title style, fall back to Heading 1 / Heading1
    if title:
        _docx_add_paragraph(doc, title, "Title", "Heading 1", "Heading1")

    # Pre-pass: isolate fenced code blocks (\x01 separators) and pipe tables.
    normalized = _normalize_md_blocks(content or "")

    for block in normalized.split("\n\n"):
        block = block.strip()
        if not block:
            continue

        # Fenced code block — recognized by the \x01 sentinel from the pre-pass.
        if "\x01" in block:
            code_lines = block.split("\x01")
            # Drop the opening/closing ``` markers
            if code_lines and code_lines[0].strip().startswith("```"):
                code_lines = code_lines[1:]
            if code_lines and code_lines[-1].strip().startswith("```"):
                code_lines = code_lines[:-1]
            _add_docx_code_block(doc, "\n".join(code_lines))
            continue

        # Horizontal rule (---, ___, ***)
        if re.fullmatch(r"[-_*]{3,}", block):
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            pBdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), "999999")
            pBdr.append(bottom)
            pPr.append(pBdr)
            continue

        # Pipe table?
        if _is_pipe_table(block):
            _add_docx_table(doc, block)
            continue

        # Multi-line bullet list?
        lines = block.splitlines()
        if all(re.match(r"^[-*]\s+", ln.strip()) for ln in lines if ln.strip()):
            list_style = _docx_resolve_style(doc, "List Paragraph", "List Bullet", "Normal")
            for ln in lines:
                ln = ln.strip().lstrip("-*").strip()
                if not ln:
                    continue
                if list_style:
                    try:
                        p = doc.add_paragraph(style=list_style)
                    except KeyError:
                        p = doc.add_paragraph()
                else:
                    p = doc.add_paragraph()
                _add_inline_runs(p, ln)
            continue

        # Multi-line numbered list?
        if all(re.match(r"^\d+\.\s+", ln.strip()) for ln in lines if ln.strip()):
            num_style = _docx_resolve_style(doc, "List Number")
            # Word auto-numbers a paragraph only when the resolved style is a
            # proper numbered-list style. If the template lacks "List Number"
            # we keep the "1. " / "2. " prefix in the text so the user doesn't
            # lose the numbering entirely.
            keep_prefix = num_style is None
            fallback_style = _docx_resolve_style(doc, "List Paragraph", "Normal")
            for ln in lines:
                stripped = re.sub(r"^\s*\d+\.\s+", "", ln) if not keep_prefix else ln.strip()
                if not stripped:
                    continue
                use_style = num_style or fallback_style
                if use_style:
                    try:
                        p = doc.add_paragraph(style=use_style)
                    except KeyError:
                        p = doc.add_paragraph()
                else:
                    p = doc.add_paragraph()
                _add_inline_runs(p, stripped)
            continue

        # Heading?
        m = re.match(r"^(#{1,6})\s+(.*)$", block)
        if m:
            level = min(len(m.group(1)), 6)
            heading_style = _docx_resolve_style(
                doc, f"Heading {level}", f"Heading{level}"
            )
            text = m.group(2).strip()
            if heading_style:
                try:
                    p = doc.add_paragraph(style=heading_style)
                    _add_inline_runs(p, text)
                    continue
                except KeyError:
                    pass
            # Last-resort fallback to python-docx built-in add_heading
            p = doc.add_heading("", level=min(level, 4))
            _add_inline_runs(p, text)
            continue

        # Blockquote (lines starting with >)
        if all(ln.strip().startswith(">") for ln in lines if ln.strip()):
            quote_style = _docx_resolve_style(doc, "Quote", "Intense Quote", "Body Text", "Normal")
            for ln in lines:
                stripped = re.sub(r"^\s*>\s?", "", ln)
                if quote_style:
                    try:
                        p = doc.add_paragraph(style=quote_style)
                    except KeyError:
                        p = doc.add_paragraph()
                else:
                    p = doc.add_paragraph()
                try:
                    p.paragraph_format.left_indent = Pt(18)
                except Exception:
                    pass
                _add_inline_runs(p, stripped)
                if not (quote_style and "uote" in quote_style):
                    for run in p.runs:
                        run.italic = True
            continue

        # Regular paragraph (preserve internal newlines as soft breaks)
        body_style = _docx_resolve_style(doc, "Body", "Body Text", "Normal")
        if body_style:
            try:
                p = doc.add_paragraph(style=body_style)
            except KeyError:
                p = doc.add_paragraph()
        else:
            p = doc.add_paragraph()
        first = True
        for ln in block.split("\n"):
            if not first:
                p.add_run().add_break()
            _add_inline_runs(p, ln)
            first = False

    doc.save(out_path)


def _make_xlsx(out_path: str, sheets: list[dict]) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    # openpyxl creates one default sheet; remove it and re-add named sheets
    default = wb.active
    wb.remove(default)

    bold = Font(bold=True)

    for i, sheet in enumerate(sheets):
        name = (sheet.get("name") or f"Sheet{i + 1}")[:31]  # Excel limit
        ws = wb.create_sheet(title=name)

        headers = sheet.get("headers") or []
        rows = sheet.get("rows") or []

        if headers:
            ws.append([str(h) for h in headers])
            for col_idx in range(1, len(headers) + 1):
                ws.cell(row=1, column=col_idx).font = bold

        for row in rows:
            if not isinstance(row, list):
                row = [row]
            ws.append(row)

        # Auto-size columns (approximate)
        max_cols = max(
            (len(headers) if headers else 0),
            max((len(r) for r in rows if isinstance(r, list)), default=0),
        )
        for col_idx in range(1, max_cols + 1):
            col_letter = ws.cell(row=1, column=col_idx).column_letter
            max_len = 10
            for cell in ws[col_letter]:
                v = cell.value
                if v is not None:
                    max_len = max(max_len, min(len(str(v)) + 2, 50))
            ws.column_dimensions[col_letter].width = max_len

    # If no sheets were added (defensive), add an empty one
    if not wb.sheetnames:
        wb.create_sheet(title="Sheet1")

    wb.save(out_path)


def _pick_pptx_layout(idx: int, total: int, slide_data: dict) -> str:
    """Auto-pick the best demo layout for a slide based on its content shape."""
    # Cover always for first slide
    if idx == 0:
        return "cover"
    # Closing: last slide with thank-you-like title
    if idx == total - 1:
        title = (slide_data.get("title") or "").lower()
        if any(k in title for k in ("terima kasih", "thank", "q&a", "questions", "tanya", "closing")):
            return "closing"
    # Structured fields drive layout choice
    if "stats" in slide_data:
        return "stat_callouts"
    if "phases" in slide_data or "timeline" in slide_data:
        return "timeline"
    if "left" in slide_data and "right" in slide_data:
        l, r = slide_data["left"], slide_data["right"]
        # Lists → comparison. Dicts with bullets/features → comparison.
        # Plain strings → two_column.
        def _looks_like_comparison(side):
            if isinstance(side, list):
                return True
            if isinstance(side, dict):
                return bool(side.get("bullets") or side.get("features") or side.get("name") or side.get("title"))
            return False
        if _looks_like_comparison(l) or _looks_like_comparison(r):
            return "comparison"
        return "two_column"
    if "items" in slide_data:
        items = slide_data["items"]
        if isinstance(items, list):
            if len(items) == 4:
                return "grid_2x2"
            if len(items) >= 1:
                return "icon_rows"
    # Section divider: only title + subtitle, no body
    if slide_data.get("subtitle") and not (slide_data.get("bullets") or slide_data.get("content")):
        return "section"
    # Fallback: two-column with text on left, image placeholder on right
    return "two_column"


def _duplicate_slide(prs, source_slide):
    """Duplicate a slide and append to the end of the presentation.

    Standard python-pptx workaround: add_slide with same layout, clear
    auto-generated shapes, deep-copy each shape's XML from the source.
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    # Remove default shapes that add_slide created from layout
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    # Copy each shape from the source slide
    for shape in source_slide.shapes:
        new_shape = deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_shape)
    return new_slide


def _remove_slide(prs, index: int) -> None:
    """Remove the slide at the given index."""
    sld_id_lst = prs.slides._sldIdLst
    sld_id = list(sld_id_lst)[index]
    # Drop the relationship from the presentation part
    rId = sld_id.attrib[
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    ]
    prs.part.drop_rel(rId)
    sld_id_lst.remove(sld_id)


_MD_MARKER_RE = re.compile(r"\*\*|`[^`\n]+`|~~|\[[^\]\n]+\]\([^)\n]+\)|(?<![\w*])\*[^*\n]+\*(?![\w*])|(?<![\w_])_[^_\n]+_(?![\w_])")


def _set_shape_text(shape, text: str) -> None:
    """Replace a shape's text while preserving font formatting from its first run.

    Removes additional paragraphs/runs, keeps the first one as the style anchor.
    Inline markdown (**bold**, *italic*, `code`, ~~strike~~, [label](url)) is
    rendered as multiple runs so it doesn't leak literal markers into the slide.
    Arabic text gets Arial font automatically.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    text = str(text)

    if not tf.paragraphs:
        tf.text = text
        return

    first_p = tf.paragraphs[0]
    # Remove other paragraphs (their XML)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)

    # Fast path: no markdown markers — preserve original behavior exactly
    if not _MD_MARKER_RE.search(text):
        if first_p.runs:
            first_run = first_p.runs[0]
            for r in list(first_p.runs[1:]):
                r._r.getparent().remove(r._r)
            first_run.text = text
            if _has_arabic(text):
                first_run.font.name = "Arial"
        else:
            first_p.text = text
        return

    # Markdown path: tokenize and emit one run per styled segment.
    tokens = _tokenize_inline_md(text)

    # Capture anchor font from the first existing run (if any) so we can copy
    # its font name/size onto each new run. Then clear all existing runs.
    anchor_name = None
    anchor_size = None
    if first_p.runs:
        anchor = first_p.runs[0]
        try:
            anchor_name = anchor.font.name
        except Exception:
            anchor_name = None
        try:
            anchor_size = anchor.font.size
        except Exception:
            anchor_size = None
        for r in list(first_p.runs):
            r._r.getparent().remove(r._r)

    for kind, payload in tokens:
        if kind == "link":
            label, url = payload  # type: ignore[misc]
            run = first_p.add_run()
            run.text = str(label)
            try:
                run.hyperlink.address = str(url)
            except Exception:
                # Fallback: append URL in parentheses so the link is visible
                run.text = f"{label} ({url})"
            seg_text = str(label)
        else:
            seg_text = str(payload)
            run = first_p.add_run()
            run.text = seg_text
            if kind == "bold":
                run.font.bold = True
            elif kind == "italic":
                run.font.italic = True
            elif kind == "code":
                run.font.name = "Consolas"

        # Apply anchor font/size unless this run set its own (code uses Consolas)
        if kind != "code" and anchor_name:
            try:
                run.font.name = anchor_name
            except Exception:
                pass
        if anchor_size:
            try:
                run.font.size = anchor_size
            except Exception:
                pass
        if _has_arabic(seg_text):
            run.font.name = "Arial"


def _pptx_apply_para_align(paragraph, alignment: str) -> None:
    """Map our 'left'|'center'|'right' to python-pptx PP_ALIGN."""
    try:
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return
    if alignment == "center":
        paragraph.alignment = PP_ALIGN.CENTER
    elif alignment == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    else:
        paragraph.alignment = PP_ALIGN.LEFT


def _pptx_set_cell_fill(cell, hex_color: str) -> None:
    """Solid-color background for a python-pptx table cell."""
    try:
        from pptx.dml.color import RGBColor
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    except Exception:
        pass


def _pptx_render_table(slide, table_block: str, left, top, width, height) -> None:
    """Add a native PowerPoint table for a markdown pipe-table at the given rect."""
    lines = [ln for ln in table_block.splitlines() if ln.strip()]
    if not lines:
        return
    header = _split_pipe_row(lines[0])
    aligns = _parse_table_alignments(lines[1]) if len(lines) > 1 else []
    body_rows = [_split_pipe_row(ln) for ln in lines[2:]]
    ncols = max([len(header)] + [len(r) for r in body_rows]) if body_rows else len(header)
    nrows = 1 + len(body_rows)
    while len(aligns) < ncols:
        aligns.append("left")

    try:
        tbl_shape = slide.shapes.add_table(nrows, max(1, ncols), left, top, width, height)
    except Exception as e:
        logger.warning(f"pptx add_table failed: {e}")
        return
    tbl = tbl_shape.table

    # Header row
    for ci in range(ncols):
        text = header[ci] if ci < len(header) else ""
        cell = tbl.cell(0, ci)
        _set_shape_text(cell, text)
        _pptx_set_cell_fill(cell, "D9D9D9")
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
            _pptx_apply_para_align(para, aligns[ci])
    # Body rows
    for ri, row in enumerate(body_rows, start=1):
        for ci in range(ncols):
            text = row[ci] if ci < len(row) else ""
            cell = tbl.cell(ri, ci)
            _set_shape_text(cell, text)
            for para in cell.text_frame.paragraphs:
                _pptx_apply_para_align(para, aligns[ci])


def _pptx_extract_table_from_text(text: str) -> tuple[str, str] | None:
    """If `text` is exactly a pipe-table (optionally with blank-line padding),
    return (table_block, ''). If a table is followed/preceded by non-table text,
    return None so the caller can keep rendering as text."""
    lines = (text or "").splitlines()
    # Strip leading/trailing blank lines
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    if len(lines) < 2:
        return None
    if not lines[0].strip().startswith("|"):
        return None
    if not _is_pipe_separator_line(lines[1].strip()):
        return None
    # Every remaining line must be a pipe row
    if not all(ln.strip().startswith("|") for ln in lines if ln.strip()):
        return None
    return "\n".join(lines), ""


def _pptx_render_text(slide, shape, text: str) -> None:
    """Wrap _set_shape_text with pipe-table detection. If the text is a pure
    pipe-table block, replace the placeholder shape with a table at the same
    position; otherwise render text normally."""
    extracted = _pptx_extract_table_from_text(text or "")
    if extracted is None:
        _set_shape_text(shape, text or "")
        return
    table_block, _ = extracted
    # Capture shape rect, then remove the placeholder and emit table.
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except Exception:
        # Shape has no usable position — fall back to plain text rendering.
        _set_shape_text(shape, text or "")
        return
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except Exception:
        pass
    _pptx_render_table(slide, table_block, left, top, width, height)


def _populate_cover(slide, data: dict, doc_title: str) -> None:
    # shape[1]=title, [2]=subtitle, [3]=meta
    title = data.get("title") or doc_title or ""
    subtitle = data.get("subtitle") or data.get("content") or ""
    meta = data.get("meta") or ""
    shapes = list(slide.shapes)
    if len(shapes) > 1 and title:
        _set_shape_text(shapes[1], title)
    if len(shapes) > 2:
        _set_shape_text(shapes[2], subtitle)
    if len(shapes) > 3 and meta:
        _set_shape_text(shapes[3], meta)


def _populate_section(slide, data: dict, section_no: int) -> None:
    # shape[1]=number, [2]=title, [3]=description
    shapes = list(slide.shapes)
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("number") or f"{section_no:02d}")
    if len(shapes) > 2:
        _set_shape_text(shapes[2], data.get("title") or "")
    if len(shapes) > 3:
        _set_shape_text(shapes[3], data.get("subtitle") or data.get("content") or "")


def _populate_two_column(slide, data: dict) -> None:
    # shape[0]=title, [1]=subtitle, [2]=left text, [4]=right placeholder
    shapes = list(slide.shapes)
    if len(shapes) > 0:
        _set_shape_text(shapes[0], data.get("title") or "")
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("subtitle") or "")
    # Left content: prefer "left" string, else bullets joined, else content
    left = data.get("left")
    if isinstance(left, list):
        left = "\n".join(f"• {x}" for x in left)
    elif not left:
        bullets = data.get("bullets") or []
        if bullets:
            left = "\n".join(f"• {b}" for b in bullets)
        else:
            left = data.get("content") or ""
    if len(shapes) > 2:
        _pptx_render_text(slide, shapes[2], left)
    # Right: if user provided right text, replace placeholder
    right = data.get("right")
    if right and len(shapes) > 4:
        if isinstance(right, list):
            right = "\n".join(f"• {x}" for x in right)
        _pptx_render_text(slide, shapes[4], right)


def _populate_icon_rows(slide, data: dict) -> None:
    # 4 item slots at indices [4,5]/[8,9]/[12,13]/[16,17] (name+desc).
    # Numbers at [3,7,11,15].
    shapes = list(slide.shapes)
    if len(shapes) > 0:
        _set_shape_text(shapes[0], data.get("title") or "")
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("subtitle") or "")
    items = data.get("items") or []
    item_count = min(len(items), 4)
    for i in range(4):
        base = 3 + i * 4
        if i < item_count:
            item = items[i]
            if isinstance(item, dict):
                name = item.get("name") or item.get("title") or ""
                desc = item.get("desc") or item.get("description") or ""
            else:
                name = str(item)
                desc = ""
        else:
            # Clear unused slot
            name = desc = ""
        if base + 1 < len(shapes):
            _set_shape_text(shapes[base + 1], name)
        if base + 2 < len(shapes):
            _set_shape_text(shapes[base + 2], desc)
        # Hide the number circle too if unused (keep its shape but blank text)
        if i >= item_count and base < len(shapes):
            _set_shape_text(shapes[base], "")


def _populate_stat_callouts(slide, data: dict) -> None:
    # 3 stat slots: big at [3,7,11], label at [5,9,13].
    shapes = list(slide.shapes)
    if len(shapes) > 0:
        _set_shape_text(shapes[0], data.get("title") or "")
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("subtitle") or "")
    stats = data.get("stats") or []
    big_positions = [3, 7, 11]
    label_positions = [5, 9, 13]
    stat_count = min(len(stats), 3)
    for i in range(3):
        if i < stat_count:
            stat = stats[i]
            if isinstance(stat, dict):
                value = stat.get("value") or stat.get("number") or ""
                label = stat.get("label") or stat.get("desc") or ""
            else:
                value = str(stat)
                label = ""
        else:
            value = label = ""
        if big_positions[i] < len(shapes):
            _set_shape_text(shapes[big_positions[i]], value)
        if label_positions[i] < len(shapes):
            _set_shape_text(shapes[label_positions[i]], label)


def _populate_grid_2x2(slide, data: dict) -> None:
    # 4 quadrants at [4,5]/[8,9]/[12,13]/[16,17].
    shapes = list(slide.shapes)
    if len(shapes) > 0:
        _set_shape_text(shapes[0], data.get("title") or "")
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("subtitle") or "")
    items = data.get("items") or []
    quadrant_starts = [4, 8, 12, 16]
    item_count = min(len(items), 4)
    for i in range(4):
        if i < item_count:
            item = items[i]
            if isinstance(item, dict):
                name = item.get("name") or item.get("title") or ""
                desc = item.get("desc") or item.get("description") or ""
            else:
                name = str(item)
                desc = ""
        else:
            name = desc = ""
        base = quadrant_starts[i]
        if base < len(shapes):
            _set_shape_text(shapes[base], name)
        if base + 1 < len(shapes):
            _set_shape_text(shapes[base + 1], desc)


def _populate_timeline(slide, data: dict) -> None:
    # 5 phase slots: number at [4,8,12,16,20], name at [5,9,...], desc at [6,10,...]
    shapes = list(slide.shapes)
    if len(shapes) > 0:
        _set_shape_text(shapes[0], data.get("title") or "")
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("subtitle") or "")
    phases = data.get("phases") or data.get("timeline") or []
    starts = [4, 8, 12, 16, 20]
    phase_count = min(len(phases), 5)
    for i in range(5):
        if i < phase_count:
            phase = phases[i]
            if isinstance(phase, dict):
                name = phase.get("name") or phase.get("title") or ""
                desc = phase.get("desc") or phase.get("description") or ""
            else:
                name = str(phase)
                desc = ""
        else:
            name = desc = ""
        base = starts[i]
        if base + 1 < len(shapes):
            _set_shape_text(shapes[base + 1], name)
        if base + 2 < len(shapes):
            _set_shape_text(shapes[base + 2], desc)
        # Hide number if unused
        if i >= phase_count and base < len(shapes):
            _set_shape_text(shapes[base], "")


def _populate_comparison(slide, data: dict) -> None:
    # shapes: [0]=title [1]=subtitle
    # left column: [3]=name, [4-7]=4 bullets
    # right column: [10]=REKOMENDASI badge, [11]=name, [12-15]=4 bullets
    shapes = list(slide.shapes)
    if len(shapes) > 0:
        _set_shape_text(shapes[0], data.get("title") or "")
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("subtitle") or "")
    left = data.get("left") or {}
    right = data.get("right") or {}
    # Left
    if isinstance(left, dict):
        l_name = left.get("name") or left.get("title") or "Opsi A"
        l_bullets = left.get("bullets") or left.get("features") or []
    elif isinstance(left, list):
        l_name = "Opsi A"
        l_bullets = left
    else:
        l_name = str(left)
        l_bullets = []
    if len(shapes) > 3:
        _set_shape_text(shapes[3], l_name)
    for i in range(4):
        if 4 + i < len(shapes):
            txt = f"•  {l_bullets[i]}" if i < len(l_bullets) else ""
            _set_shape_text(shapes[4 + i], txt)
    # Right
    if isinstance(right, dict):
        r_name = right.get("name") or right.get("title") or "Opsi B"
        r_bullets = right.get("bullets") or right.get("features") or []
        r_recommended = right.get("recommended", True)
    elif isinstance(right, list):
        r_name = "Opsi B"
        r_bullets = right
        r_recommended = True
    else:
        r_name = str(right)
        r_bullets = []
        r_recommended = False
    if len(shapes) > 10 and not r_recommended:
        _set_shape_text(shapes[10], "")
    if len(shapes) > 11:
        _set_shape_text(shapes[11], r_name)
    for i in range(4):
        if 12 + i < len(shapes):
            txt = f"✓  {r_bullets[i]}" if i < len(r_bullets) else ""
            _set_shape_text(shapes[12 + i], txt)


def _populate_closing(slide, data: dict) -> None:
    # shapes: [0]=title, [1]=subtitle, [3]=contact
    shapes = list(slide.shapes)
    if len(shapes) > 0:
        _set_shape_text(shapes[0], data.get("title") or "Terima Kasih")
    if len(shapes) > 1:
        _set_shape_text(shapes[1], data.get("subtitle") or "Diskusi dan Pertanyaan")
    if len(shapes) > 3:
        _set_shape_text(shapes[3], data.get("meta") or data.get("contact") or "")


_POPULATORS = {
    "cover": _populate_cover,
    "section": lambda slide, data, idx=0: _populate_section(slide, data, idx),
    "two_column": _populate_two_column,
    "icon_rows": _populate_icon_rows,
    "stat_callouts": _populate_stat_callouts,
    "grid_2x2": _populate_grid_2x2,
    "timeline": _populate_timeline,
    "comparison": _populate_comparison,
    "closing": _populate_closing,
}


def _make_pptx(out_path: str, title: str, slides: list[dict]) -> None:
    from pptx import Presentation

    template = _pptx_template_path()
    if template:
        prs = Presentation(template)
        # Demo slides (0-8) act as a layout library. Plan layouts BEFORE
        # mutating the deck so all references stay valid.
        demo_count = len(prs.slides)
        demo_lookup = {name: i for i, name in enumerate(_PPTX_LAYOUTS)}

        # Pick layout for each user slide first
        chosen: list[tuple[int, dict]] = []  # (demo_index, slide_data)
        section_counter = 0
        for i, slide_data in enumerate(slides):
            layout_name = _pick_pptx_layout(i, len(slides), slide_data)
            demo_idx = demo_lookup.get(layout_name, demo_lookup["two_column"])
            chosen.append((demo_idx, slide_data))

        # Clone the right demo for each user slide, then populate
        cloned_slides = []
        for demo_idx, _data in chosen:
            new_slide = _duplicate_slide(prs, prs.slides[demo_idx])
            cloned_slides.append(new_slide)

        for clone, (demo_idx, slide_data) in zip(cloned_slides, chosen):
            layout_name = _PPTX_LAYOUTS[demo_idx]
            populator = _POPULATORS.get(layout_name, _populate_two_column)
            if layout_name == "cover":
                populator(clone, slide_data, title)
            elif layout_name == "section":
                section_counter += 1
                populator(clone, slide_data, section_counter)
            else:
                populator(clone, slide_data)

        # Remove the original demo slides (first demo_count slides)
        for _ in range(demo_count):
            _remove_slide(prs, 0)

    else:
        # No template — fall back to plain layout (legacy behavior)
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        bullet_layout = prs.slide_layouts[1]
        for i, slide_data in enumerate(slides):
            s_title = (slide_data.get("title") or "").strip()
            bullets = slide_data.get("bullets") or []
            content = (slide_data.get("content") or "").strip()
            if i == 0 and title and not s_title:
                slide = prs.slides.add_slide(title_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = title
                try:
                    sub = slide.placeholders[1]
                    sub.text = content or " ".join(str(b) for b in bullets[:3])
                except (KeyError, IndexError):
                    pass
                continue
            slide = prs.slides.add_slide(bullet_layout)
            if slide.shapes.title:
                slide.shapes.title.text = s_title or f"Slide {i + 1}"
            try:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.word_wrap = True
                if bullets:
                    for j, b in enumerate(bullets):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = str(b)
                elif content:
                    tf.text = content
            except (KeyError, IndexError):
                pass

    prs.save(out_path)
